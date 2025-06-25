from pptx import Presentation
import re
import os
from PIL import Image
import io

try:
    # PowerPoint 파일 열기
    ppt_file = 'picture_2.pptx'
    prs = Presentation(ppt_file)
    
    # 이미지 저장 폴더 생성
    image_folder = 'student_images'
    if not os.path.exists(image_folder):
        os.makedirs(image_folder)
        print(f"폴더 생성: {image_folder}")
    
    print("=== PowerPoint에서 이미지 추출 중... ===")
    
    all_students = []
    extracted_images = []
    
    # 각 슬라이드 처리
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"\n슬라이드 {slide_idx} 처리 중...")
        
        # 학생 정보 추출
        student_info_list = []
        for shape in slide.shapes:
            if shape.shape_type == 19:  # Table type
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        if cell_text and "학년" in cell_text and "반" in cell_text and "번" in cell_text:
                            # 학생 정보 추출 (예: "3학년 1반 2번 권민지")
                            match = re.search(r'(\d)학년\s*(\d)반\s*(\d+)번\s*(.+)', cell_text)
                            if match:
                                grade, class_num, student_num, name = match.groups()
                                # 학번 생성: 3AABB 형식
                                student_id = f"3{int(class_num):02d}{int(student_num):02d}"
                                student_info = {
                                    '학번': student_id,
                                    '반': class_num,
                                    '번호': student_num,
                                    '이름': name.strip(),
                                    '위치': (row_idx, col_idx)
                                }
                                student_info_list.append(student_info)
                                all_students.append(student_info)
        
        # 이미지 추출
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture type
                try:
                    # 이미지 데이터 추출
                    image_data = shape.image.blob
                    image = Image.open(io.BytesIO(image_data))
                    
                    # 이미지 순서에 맞는 학생 정보 찾기
                    if image_count < len(student_info_list):
                        student = student_info_list[image_count]
                        filename = f"{student['학번']}_{student['이름']}.jpg"
                        filepath = os.path.join(image_folder, filename)
                        
                        # 이미지 저장
                        image.save(filepath, 'JPEG')
                        extracted_images.append({
                            '학번': student['학번'],
                            '이름': student['이름'],
                            '파일명': filename,
                            '파일경로': filepath
                        })
                        print(f"  저장: {filename}")
                    else:
                        # 학생 정보가 없는 경우 일반 파일명으로 저장
                        filename = f"slide{slide_idx}_image{image_count+1}.jpg"
                        filepath = os.path.join(image_folder, filename)
                        image.save(filepath, 'JPEG')
                        print(f"  저장: {filename} (학생 정보 없음)")
                    
                    image_count += 1
                    
                except Exception as e:
                    print(f"  이미지 추출 오류: {e}")
    
    print(f"\n=== 추출 완료 ===")
    print(f"총 학생 수: {len(all_students)}명")
    print(f"추출된 이미지: {len(extracted_images)}개")
    
    # 추출 결과를 CSV로 저장
    if extracted_images:
        import pandas as pd
        df_images = pd.DataFrame(extracted_images)
        df_images.to_csv('extracted_images.csv', index=False, encoding='utf-8-sig')
        print(f"추출 정보가 'extracted_images.csv'에 저장되었습니다.")
        
        print("\n=== 추출된 이미지 목록 ===")
        for img_info in extracted_images:
            print(f"{img_info['학번']} {img_info['이름']} -> {img_info['파일명']}")
    
    # 누락된 학생 확인
    print(f"\n=== 이미지 추출 현황 ===")
    print(f"학생 명단: {len(all_students)}명")
    print(f"추출 이미지: {len(extracted_images)}개")
    
    if len(all_students) != len(extracted_images):
        print("⚠️ 일부 학생의 이미지가 누락되었을 수 있습니다.")
        print("PowerPoint의 이미지와 텍스트 순서를 확인해주세요.")

except Exception as e:
    print(f"오류가 발생했습니다: {e}")
    import traceback
    traceback.print_exc()
