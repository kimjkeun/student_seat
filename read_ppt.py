from pptx import Presentation
import os

try:
    # PowerPoint 파일 열기
    ppt_file = 'picture_2.pptx'
    prs = Presentation(ppt_file)
    
    print("=== PowerPoint 파일 정보 ===")
    print(f"파일명: {ppt_file}")
    print(f"슬라이드 수: {len(prs.slides)}개")
    print()
    
    # 각 슬라이드 내용 확인
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"=== 슬라이드 {slide_idx} ===")
        
        # 텍스트 내용 추출
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
        
        if text_content:
            print("텍스트 내용:")
            for i, text in enumerate(text_content, 1):
                print(f"  {i}. {text}")
        else:
            print("텍스트 내용 없음")
        
        # 이미지 확인
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture type
                image_count += 1
        
        if image_count > 0:
            print(f"이미지 개수: {image_count}개")
        
        # 표(테이블) 확인
        table_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 19:  # Table type
                table_count += 1
                print(f"표 발견 - 행: {len(shape.table.rows)}, 열: {len(shape.table.columns)}")
                
                # 표 내용 출력
                print("표 내용:")
                for row_idx, row in enumerate(shape.table.rows):
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    print(f"  행 {row_idx+1}: {' | '.join(row_data)}")
        
        print()

except ImportError:
    print("python-pptx 라이브러리가 설치되어 있지 않습니다.")
    print("다음 명령어로 설치하세요: pip install python-pptx")
except FileNotFoundError:
    print(f"파일을 찾을 수 없습니다: {ppt_file}")
except Exception as e:
    print(f"오류가 발생했습니다: {e}")
