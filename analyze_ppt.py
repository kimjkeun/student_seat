from pptx import Presentation
import re

try:
    # PowerPoint 파일 열기
    ppt_file = 'picture_2.pptx'
    prs = Presentation(ppt_file)
    
    print("=== PowerPoint 파일 정보 ===")
    print(f"파일명: {ppt_file}")
    print(f"슬라이드 수: {len(prs.slides)}개")
    print()
    
    all_students = []
    
    # 각 슬라이드 내용 확인
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"=== 슬라이드 {slide_idx} ===")
        
        # 이미지 개수 확인
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture type
                image_count += 1
        print(f"이미지 개수: {image_count}개")
        
        # 표(테이블) 분석
        table_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 19:  # Table type
                table_count += 1
                print(f"\n표 {table_count} - 행: {len(shape.table.rows)}, 열: {len(shape.table.columns)}")
                
                # 표 내용 상세 분석
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        if cell_text and "학년" in cell_text and "반" in cell_text and "번" in cell_text:
                            # 학생 정보 추출 (예: "3학년 1반 2번 권민지")
                            match = re.search(r'(\d)학년\s*(\d)반\s*(\d+)번\s*(.+)', cell_text)
                            if match:
                                grade, class_num, student_num, name = match.groups()
                                student_info = {
                                    '학년': grade,
                                    '반': class_num,
                                    '번호': student_num,
                                    '이름': name.strip(),
                                    '슬라이드': slide_idx,
                                    '위치': f'행{row_idx+1}_열{col_idx+1}'
                                }
                                all_students.append(student_info)
                                print(f"  학생 발견: {grade}학년 {class_num}반 {student_num}번 {name.strip()}")
        
        print()
    
    # 전체 학생 명단 정리
    print("=== 전체 학생 명단 ===")
    print(f"총 학생 수: {len(all_students)}명")
    
    if all_students:
        # 반별로 정렬
        all_students.sort(key=lambda x: (int(x['반']), int(x['번호'])))
        
        print("\n=== 반별 학생 명단 ===")
        current_class = None
        for student in all_students:
            if current_class != student['반']:
                current_class = student['반']
                print(f"\n{student['반']}반:")
            print(f"  {student['번호']:2s}번 {student['이름']}")
        
        # CSV로 저장
        import pandas as pd
        df_students = pd.DataFrame(all_students)
        df_students.to_csv('student_list.csv', index=False, encoding='utf-8-sig')
        print(f"\n학생 명단이 'student_list.csv' 파일로 저장되었습니다.")
        
        # 반별 통계
        print("\n=== 반별 학생 수 ===")
        class_counts = df_students['반'].value_counts().sort_index()
        for class_num, count in class_counts.items():
            print(f"{class_num}반: {count}명")

except Exception as e:
    print(f"오류가 발생했습니다: {e}")
