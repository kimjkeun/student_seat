from pptx import Presentation
import re
import os
from PIL import Image
import io

def extract_images_from_ppt(ppt_path, class_num):
    """PowerPoint 파일에서 이미지를 추출하고 클래스별 폴더에 저장"""
    try:
        prs = Presentation(ppt_path)
        
        # 클래스별 폴더 생성
        class_folder = f'class_{class_num}'
        if not os.path.exists(class_folder):
            os.makedirs(class_folder)
            print(f"폴더 생성: {class_folder}")
        
        all_students = []
        extracted_images = []
        
        print(f"\n=== 클래스 {class_num} - {os.path.basename(ppt_path)} 처리 중... ===")
        
        # 각 슬라이드 처리
        for slide_idx, slide in enumerate(prs.slides, 1):
            print(f"  슬라이드 {slide_idx} 처리 중...")
            
            # 학생 정보 추출
            student_info_list = []
            for shape in slide.shapes:
                if shape.shape_type == 19:  # Table type
                    for row_idx, row in enumerate(shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()
                            if cell_text and "학년" in cell_text and "반" in cell_text and "번" in cell_text:
                                # 학생 정보 추출 (예: "3학년 1반 2번 권민지")
                                match = re.search(r'(\d)학년\s*(\d)반\s*(\d+)번\s*(.+)', cell_text)
                                if match:
                                    grade, class_num_in_text, student_num, name = match.groups()
                                    # 학번 생성: 3AABB 형식
                                    student_id = f"3{int(class_num_in_text):02d}{int(student_num):02d}"
                                    student_info = {
                                        '학번': student_id,
                                        '반': class_num_in_text,
                                        '번호': student_num,
                                        '이름': name.strip(),
                                        '위치': (row_idx, col_idx)
                                    }
                                    student_info_list.append(student_info)
                                    all_students.append(student_info)
            
            # 이미지 추출
            image_count = 0
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture type
                    try:
                        # 이미지 데이터 추출
                        image_data = shape.image.blob
                        image = Image.open(io.BytesIO(image_data))
                        
                        # 이미지 순서에 맞는 학생 정보 찾기
                        if image_count < len(student_info_list):
                            student = student_info_list[image_count]
                            filename = f"{student['학번']}_{student['이름']}.jpg"
                            filepath = os.path.join(class_folder, filename)
                            
                            # 이미지 저장
                            image.save(filepath, 'JPEG')
                            extracted_images.append({
                                '클래스': class_num,
                                '학번': student['학번'],
                                '이름': student['이름'],
                                '파일명': filename,
                                '파일경로': filepath
                            })
                            print(f"    저장: {filename}")
                        else:
                            # 학생 정보가 없는 경우 일반 파일명으로 저장
                            filename = f"class{class_num}_slide{slide_idx}_image{image_count+1}.jpg"
                            filepath = os.path.join(class_folder, filename)
                            image.save(filepath, 'JPEG')
                            print(f"    저장: {filename} (학생 정보 없음)")
                        
                        image_count += 1
                        
                    except Exception as e:
                        print(f"    이미지 추출 오류: {e}")
        
        print(f"  클래스 {class_num}: {len(all_students)}명, {len(extracted_images)}개 이미지 추출")
        return all_students, extracted_images
        
    except Exception as e:
        print(f"클래스 {class_num} 처리 오류: {e}")
        return [], []

# 메인 실행
try:
    image_folder = 'c:/Users/user/Documents/00_temp/data/image'
    os.chdir(image_folder)
    
    print("=== 클래스별 학생 이미지 추출 시작 ===")
    
    # PowerPoint 파일 목록 확인
    ppt_files = [f for f in os.listdir('.') if f.endswith('.pptx')]
    print(f"발견된 PowerPoint 파일: {len(ppt_files)}개")
    
    all_class_students = []
    all_class_images = []
    
    # 각 PowerPoint 파일 처리
    for ppt_file in sorted(ppt_files):
        # 파일명에서 클래스 번호 추출 (picture_n.pptx에서 n)
        match = re.search(r'picture_(\d+)\.pptx', ppt_file)
        if match:
            class_num = match.group(1)
            students, images = extract_images_from_ppt(ppt_file, class_num)
            all_class_students.extend(students)
            all_class_images.extend(images)
        else:
            print(f"클래스 번호를 찾을 수 없음: {ppt_file}")
    
    print(f"\n=== 전체 추출 완료 ===")
    print(f"총 학생 수: {len(all_class_students)}명")
    print(f"총 이미지 수: {len(all_class_images)}개")
    
    # 결과를 CSV로 저장
    if all_class_images:
        import pandas as pd
        df_all_images = pd.DataFrame(all_class_images)
        df_all_images.to_csv('all_class_images.csv', index=False, encoding='utf-8-sig')
        print(f"전체 추출 정보가 'all_class_images.csv'에 저장되었습니다.")
        
        # 클래스별 통계
        print("\n=== 클래스별 학생 수 ===")
        class_counts = df_all_images['클래스'].value_counts().sort_index()
        for class_num, count in class_counts.items():
            print(f"클래스 {class_num}: {count}명")
    
    # 생성된 폴더 확인
    print(f"\n=== 생성된 클래스 폴더 ===")
    class_folders = [d for d in os.listdir('.') if d.startswith('class_') and os.path.isdir(d)]
    for folder in sorted(class_folders):
        file_count = len([f for f in os.listdir(folder) if f.endswith('.jpg')])
        print(f"{folder}: {file_count}개 이미지")

except Exception as e:
    print(f"전체 처리 오류: {e}")
    import traceback
    traceback.print_exc()
